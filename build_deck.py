"""
Generates the SyncFlow Automator pitch deck (SyncFlow_Automator_Deck.pptx)
using python-pptx. Theme mirrors the app's own dark UI.

Run:  python build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Palette (matches main.py) ----------
BG        = RGBColor(0x12, 0x12, 0x14)   # app background
PANEL     = RGBColor(0x1E, 0x1E, 0x24)   # entry / panel
PANEL2    = RGBColor(0x2D, 0x2D, 0x34)   # button grey
BLUE      = RGBColor(0x00, 0x7A, 0xCC)   # primary accent
BLUE_LT   = RGBColor(0x33, 0xA1, 0xFF)
GREEN     = RGBColor(0x28, 0xA7, 0x45)   # rename action
AMBER     = RGBColor(0xF5, 0xA6, 0x23)
RED       = RGBColor(0xE0, 0x5A, 0x5A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0xB0, 0xB0, 0xB8)
GREY_DIM  = RGBColor(0x8A, 0x8A, 0x92)

HERE = os.path.dirname(os.path.abspath(__file__))
SHOT_MAIN = os.path.join(HERE, "screenshots", "main_interface.png")
SHOT_OVER = os.path.join(HERE, "screenshots", "overlay_popup.png")

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Segoe UI"
MONO = "Consolas"


# ---------- helpers ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    # send background to back
    r._element.addprevious(r._element)  # no-op safeguard
    return s

def _no_autosize(tf):
    # prevent python-pptx from resizing text frames
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shadow=False, round_=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    _no_autosize(tf)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, *rest) in para:
            fnt = rest[0] if rest else FONT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = fnt
    return tb

def accent_bar(s, x=0.0, y=0.0, w=13.333, h=0.14, color=BLUE):
    box(s, x, y, w, h, fill=color)

def chip(s, x, y, label, color=BLUE, tw=1.9):
    box(s, x, y, tw, 0.42, fill=color, round_=True)
    text(s, x, y+0.02, tw, 0.38, [[(label, 12, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)

def page_header(s, kicker, title, num):
    accent_bar(s)
    text(s, 0.7, 0.42, 11, 0.4, [[(kicker.upper(), 12, BLUE_LT, True)]], space_after=0)
    text(s, 0.7, 0.72, 12, 0.9, [[(title, 30, WHITE, True)]], space_after=0)
    # footer
    text(s, 0.7, 7.02, 6, 0.35, [[("SyncFlow Automator", 9, GREY_DIM, False)]], space_after=0)
    text(s, 11.4, 7.02, 1.3, 0.35, [[(f"{num:02d}", 9, GREY_DIM, False)]],
         align=PP_ALIGN.RIGHT, space_after=0)

def bullet(prefix_color):
    return prefix_color


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = slide()
accent_bar(s, h=0.2)
# big monogram tile
box(s, 0.9, 2.35, 1.5, 1.5, fill=BLUE, round_=True)
text(s, 0.9, 2.35, 1.5, 1.5, [[("⚡", 46, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)

text(s, 2.7, 2.35, 10, 1.2, [[("SyncFlow Automator", 46, WHITE, True)]], space_after=0)
text(s, 2.75, 3.5, 10, 0.7,
     [[("Zero-friction folder generation & context-aware renaming ", 18, GREY, False)],
      [("for multi-character video post-production pipelines", 18, GREY, False)]],
     space_after=2, line_spacing=1.05)

# meta chips
chip(s, 2.75, 4.75, "Python · Tkinter", color=PANEL2, tw=2.3)
chip(s, 5.2, 4.75, "Windows + macOS", color=PANEL2, tw=2.3)
chip(s, 7.65, 4.75, "Zero dependencies", color=PANEL2, tw=2.3)

text(s, 0.9, 6.35, 12, 0.5,
     [[("Product Overview & Improvement Roadmap", 14, BLUE_LT, True),
       ("     |     v1.1.0", 14, GREY_DIM, False)]], space_after=0)


# =====================================================================
# SLIDE 2 — THE PROBLEM
# =====================================================================
s = slide()
page_header(s, "The Problem", "Manual folder work costs hours — and invites errors", 2)

cards = [
    ("🗂  Deeply nested trees", "Every character × every scene needs 5 format "
     "folders. One reel can mean hundreds of folders built by hand.", RED),
    ("⏱  Dozens of man-hours", "Repetitive setup and delivery prep drains editor "
     "time that should go to the actual cut.", AMBER),
    ("⚠  Naming mistakes ship", "Hand-typed export names drift from convention — "
     "wrong shot, wrong suffix, broken sync downstream.", RED),
    ("🔁  No standard, no memory", "Each editor does it slightly differently; "
     "nothing enforces the pipeline's structure.", AMBER),
]
x = 0.7; y = 2.0; w = 5.9; h = 2.15; gapx = 0.35; gapy = 0.3
for i, (t, d, c) in enumerate(cards):
    cx = x + (i % 2) * (w + gapx)
    cy = y + (i // 2) * (h + gapy)
    box(s, cx, cy, w, h, fill=PANEL, round_=True)
    box(s, cx, cy, 0.12, h, fill=c)  # accent stripe
    text(s, cx+0.35, cy+0.28, w-0.6, 0.6, [[(t, 18, WHITE, True)]], space_after=0)
    text(s, cx+0.35, cy+0.95, w-0.6, 1.1, [[(d, 13.5, GREY, False)]],
         space_after=0, line_spacing=1.08)


# =====================================================================
# SLIDE 3 — THE SOLUTION
# =====================================================================
s = slide()
page_header(s, "The Solution", "One lightweight desktop app, three jobs", 3)

feats = [
    ("A", "Instant Folder Trees", "Deploy a full character → scene → format "
     "structure in one click.", BLUE),
    ("B", "Context-Aware Rename", "Renames the selected file to spec by reading "
     "its own folder path.", GREEN),
    ("C", "Integrated Quick Search", "Jump into any deeply-nested folder as you "
     "type — no digging.", AMBER),
]
cw = 3.9; gx = 0.35; sx = 0.7; sy = 2.15; ch = 3.4
for i, (letter, title, desc, c) in enumerate(feats):
    cx = sx + i * (cw + gx)
    box(s, cx, sy, cw, ch, fill=PANEL, round_=True)
    box(s, cx+0.4, sy+0.45, 0.95, 0.95, fill=c, round_=True)
    text(s, cx+0.4, sy+0.45, 0.95, 0.95, [[(letter, 30, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, cx+0.4, sy+1.65, cw-0.8, 0.8, [[(title, 19, WHITE, True)]], space_after=0)
    text(s, cx+0.4, sy+2.35, cw-0.8, 1.0, [[(desc, 13.5, GREY, False)]],
         space_after=0, line_spacing=1.1)

text(s, 0.7, 5.95, 12, 0.7,
     [[("Native ", 14, GREY, False), ("tkinter", 14, BLUE_LT, True, MONO),
       (" only — no pip installs, no background hooks, near-zero latency on remote "
        "cloud workstations.", 14, GREY, False)]], space_after=0)


# =====================================================================
# SLIDE 4 — FEATURE A: FOLDER TREE
# =====================================================================
s = slide()
page_header(s, "Feature A", "Directory Tree Generator — 3 modes", 4)

# left: mode list
mx = 0.7
modes = [
    ("Mode A", "All Structure", "Full macro ecosystem: Projects / Media / Exports "
     "and every delivery folder.", BLUE),
    ("Mode B", "Exports Only", "Just the delivery dirs — Raw, Source, Test, Master "
     "Timeline, Character Exports.", PANEL2),
    ("Mode C", "Character Exports", "The core engine: loops characters × scenes and "
     "builds 5 format folders in each.", GREEN),
]
my = 2.05
for tag, name, desc, c in modes:
    box(s, mx, my, 6.0, 1.45, fill=PANEL, round_=True)
    box(s, mx, my, 0.12, 1.45, fill=c)
    text(s, mx+0.35, my+0.18, 5.5, 0.5,
         [[(tag + "  ", 15, c if c != PANEL2 else BLUE_LT, True), (name, 15, WHITE, True)]],
         space_after=0)
    text(s, mx+0.35, my+0.68, 5.5, 0.7, [[(desc, 12.5, GREY, False)]],
         space_after=0, line_spacing=1.05)
    my += 1.62

# right: tree render
tx = 7.15
box(s, tx, 2.05, 5.45, 4.6, fill=PANEL, round_=True)
text(s, tx+0.35, 2.25, 5.0, 0.4, [[("Generated by Mode C", 12, GREY_DIM, True)]], space_after=0)
tree = [
    ("Character Exports/", BLUE_LT, 0),
    ("└─ bubu/", WHITE, 0),
    ("   └─ Scene_01/", WHITE, 1),
    ("      ├─ MP4/", GREEN, 2),
    ("      ├─ Quicktime/", GREEN, 2),
    ("      ├─ SyncSO/", GREEN, 2),
    ("      ├─ Lipdub/", GREEN, 2),
    ("      └─ Audio/", GREEN, 2),
    ("   └─ Scene_02/  …", GREY, 1),
    ("└─ meow/  …", GREY, 0),
]
para = [[(line, 14, color, False, MONO)] for line, color, _ in tree]
text(s, tx+0.45, 2.75, 5.0, 3.7, para, space_after=3, line_spacing=1.02)


# =====================================================================
# SLIDE 5 — FEATURE B: NAMING MATRIX
# =====================================================================
s = slide()
page_header(s, "Feature B", "Context-Aware Renaming Matrix", 5)

text(s, 0.7, 1.95, 12, 0.6,
     [[("Select a file → click Rename → type the shot number. The app reads the "
        "folder path and builds the spec name automatically.", 14, GREY, False)]],
     space_after=0, line_spacing=1.05)

# table
tbl_x, tbl_y = 0.7, 2.75
col_w = [2.2, 2.0, 1.8, 5.6]
headers = ["Format Folder", "Suffix", "Extension", "Filename Output"]
rows = [
    ("MP4", "—", ".mp4", "Show_Reel_Scene_Char_Shot_01.mp4", GREEN),
    ("Quicktime", "—", ".mov", "Show_Reel_Scene_Char_Shot_01.mov", GREEN),
    ("Audio", "_Audio", "kept*", "Show_Reel_Scene_Char_Shot_01_Audio.wav", BLUE_LT),
    ("SyncSO", "_Sync_so", "kept*", "Show_Reel_Scene_Char_Shot_01_Sync_so.mp4", AMBER),
    ("Lipdub", "_Lipdub", "kept*", "Show_Reel_Scene_Char_Shot_01_Lipdub.mov", AMBER),
]
# header row
hx = tbl_x
box(s, tbl_x, tbl_y, sum(col_w), 0.5, fill=BLUE)
for j, htext in enumerate(headers):
    text(s, hx+0.15, tbl_y+0.06, col_w[j]-0.2, 0.4, [[(htext, 12.5, WHITE, True)]],
         space_after=0, anchor=MSO_ANCHOR.MIDDLE)
    hx += col_w[j]
# body rows
ry = tbl_y + 0.5
for (fmt, suf, ext, out, c) in rows:
    box(s, tbl_x, ry, sum(col_w), 0.6, fill=PANEL)
    cells = [(fmt, WHITE, True, FONT), (suf, GREY, False, MONO),
             (ext, GREY, False, MONO), (out, c, False, MONO)]
    cxp = tbl_x
    for j, (val, col, bold, fnt) in enumerate(cells):
        text(s, cxp+0.15, ry+0.11, col_w[j]-0.2, 0.4, [[(val, 11.5, col, bold, fnt)]],
             space_after=0, anchor=MSO_ANCHOR.MIDDLE)
        cxp += col_w[j]
    ry += 0.63

text(s, 0.7, ry+0.05, 12, 0.5,
     [[("* v1.1 fix: the file's real extension is preserved — renaming never fakes a "
        "codec (a .wav stays .wav).", 12, GREEN, True)]], space_after=0)


# =====================================================================
# SLIDE 6 — FEATURE C: QUICK SEARCH
# =====================================================================
s = slide()
page_header(s, "Feature C", "Integrated Quick Search", 6)

pts = [
    ("Type-to-filter", "Live, recursive filter across the whole project tree."),
    ("Multi-word match", "“bubu Scene 1” isolates exactly that path matrix."),
    ("Double-click to open", "Jump straight into the folder in Explorer / Finder."),
    ("v1.1: instant", "Tree is cached & filtered in memory — no re-walk per keystroke."),
]
py = 2.15
for t, d in pts:
    box(s, 0.7, py, 0.16, 0.16, fill=BLUE)  # square bullet
    text(s, 1.05, py-0.12, 5.6, 0.5, [[(t + " — ", 15, WHITE, True), (d, 14, GREY, False)]],
         space_after=0, line_spacing=1.05)
    py += 1.02

# right mock search panel
qx = 7.0
box(s, qx, 2.05, 5.6, 4.5, fill=PANEL, round_=True)
box(s, qx+0.35, 2.4, 4.9, 0.5, fill=PANEL2, round_=True)
text(s, qx+0.5, 2.47, 4.6, 0.4, [[("bubu Scene 1", 13, WHITE, False, MONO)]], space_after=0)
results = [
    "…/Character Exports/bubu/Scene_01",
    "…/Character Exports/bubu/Scene_01/MP4",
    "…/Character Exports/bubu/Scene_01/Audio",
    "…/Character Exports/bubu/Scene_01/SyncSO",
]
ryy = 3.15
for i, rr in enumerate(results):
    if i == 0:
        box(s, qx+0.35, ryy-0.05, 4.9, 0.5, fill=BLUE, round_=True)
    text(s, qx+0.5, ryy, 4.7, 0.4, [[(rr, 11.5, WHITE if i == 0 else GREY, False, MONO)]],
         space_after=0)
    ryy += 0.58


# =====================================================================
# SLIDE 7 — SCREENSHOT: MAIN UI
# =====================================================================
s = slide()
page_header(s, "Interface", "Main configuration workspace", 7)
if os.path.exists(SHOT_MAIN):
    # image is 680x596; place centered-left, scale to height ~5"
    ih = Inches(5.0); iw = Inches(5.0 * 680/596)
    s.shapes.add_picture(SHOT_MAIN, Inches(0.9), Inches(1.9), height=ih)
    # caption card right
    cxx = 6.6
    box(s, cxx, 2.1, 6.0, 4.6, fill=PANEL, round_=True)
    labels = [
        ("Target directory + Browse", "Point it at any project root."),
        ("Project / Reel / Characters / Scenes", "The four inputs that drive generation."),
        ("Mode A · B · C buttons", "One click each — colour-coded by scope."),
        ("⚡ Rename Selected File", "The green context-rename trigger."),
        ("Quick Search + results", "Filter and open any nested folder."),
    ]
    lyy = 2.45
    for t, d in labels:
        box(s, cxx+0.35, lyy+0.06, 0.14, 0.14, fill=BLUE)
        text(s, cxx+0.65, lyy-0.08, 5.1, 0.5,
             [[(t, 13.5, WHITE, True), ("  —  " + d, 12.5, GREY, False)]], space_after=0)
        lyy += 0.82
else:
    text(s, 0.9, 3, 8, 1, [[("(screenshot not found)", 16, RED, False)]])


# =====================================================================
# SLIDE 8 — SCREENSHOT: OVERLAY
# =====================================================================
s = slide()
page_header(s, "Interface", "The context rename overlay", 8)
if os.path.exists(SHOT_OVER):
    s.shapes.add_picture(SHOT_OVER, Inches(0.9), Inches(1.9), height=Inches(5.0))
    cxx = 6.6
    box(s, cxx, 2.1, 6.0, 4.6, fill=PANEL, round_=True)
    text(s, cxx+0.4, 2.4, 5.2, 0.6, [[("Detected context, zero typing", 18, WHITE, True)]],
         space_after=0)
    steps = [
        "1.  You highlight an export file in Explorer / Finder.",
        "2.  App reads the path — Show, Reel, Scene, Character, Format.",
        "3.  A centered overlay shows what it detected.",
        "4.  You type only the shot number and hit Enter.",
        "5.  File is renamed to the exact delivery spec.",
    ]
    syy = 3.15
    for st in steps:
        text(s, cxx+0.4, syy, 5.3, 0.5, [[(st, 13.5, GREY, False)]], space_after=0,
             line_spacing=1.05)
        syy += 0.62
    text(s, cxx+0.4, syy+0.05, 5.3, 0.5,
         [[("No global keyboard hooks → no system lag.", 12.5, GREEN, True)]], space_after=0)
else:
    text(s, 0.9, 3, 8, 1, [[("(screenshot not found)", 16, RED, False)]])


# =====================================================================
# SLIDE 9 — HOW IT WORKS
# =====================================================================
s = slide()
page_header(s, "Under the Hood", "How the rename engine thinks", 9)

flow = [
    ("Capture", "Read selected path\n(CF_HDROP on Win,\nAppleScript on Mac)", BLUE),
    ("Anchor", "Find 'Character\nExports' in the path\ncase-insensitively", GREEN),
    ("Parse", "Walk up for Reel &\nShow; down for Char,\nScene, Format", AMBER),
    ("Build", "Compose spec name,\npreserve real ext,\nzero-pad shot", BLUE),
]
fx = 0.7; fw = 2.75; fy = 2.35; gap = 0.42
for i, (t, d, c) in enumerate(flow):
    cx = fx + i * (fw + gap)
    box(s, cx, fy, fw, 2.5, fill=PANEL, round_=True)
    box(s, cx, fy, fw, 0.6, fill=c, round_=True)
    text(s, cx, fy+0.08, fw, 0.5, [[(f"{i+1}. {t}", 16, WHITE, True)]],
         align=PP_ALIGN.CENTER, space_after=0)
    text(s, cx+0.25, fy+0.85, fw-0.5, 1.5, [[(d, 12.5, GREY, False, MONO)]],
         space_after=0, line_spacing=1.15)
    if i < len(flow) - 1:
        text(s, cx+fw-0.02, fy+0.85, 0.5, 0.5, [[("→", 22, BLUE_LT, True)]],
             align=PP_ALIGN.CENTER, space_after=0)

box(s, 0.7, 5.35, 11.9, 1.3, fill=PANEL, round_=True)
text(s, 1.0, 5.55, 11.4, 1.0,
     [[("Pure-logic core is unit-tested — ", 14, WHITE, True),
       ("parse_context_from_path()", 13, BLUE_LT, False, MONO),
       (" and ", 14, WHITE, True),
       ("build_new_name()", 13, BLUE_LT, False, MONO),
       (" run headless with no GUI.", 14, WHITE, True)],
      [("10/10 tests pass — see test_syncflow.py", 13, GREEN, True)]],
     space_after=4, line_spacing=1.1)


# =====================================================================
# SLIDE 10 — IMPACT
# =====================================================================
s = slide()
page_header(s, "Impact", "Why it matters to the team", 10)

stats = [
    ("100s → 1", "folders per reel,\nin a single click", BLUE),
    ("~0", "pip dependencies\nto install & maintain", GREEN),
    ("2 OS", "Windows 10/11\n+ native macOS", AMBER),
    ("0", "naming mistakes when\nthe path drives the name", BLUE),
]
sx = 0.7; sw = 2.95; sy = 2.2; g = 0.3
for i, (big, cap, c) in enumerate(stats):
    cx = sx + i * (sw + g)
    box(s, cx, sy, sw, 2.4, fill=PANEL, round_=True)
    text(s, cx, sy+0.4, sw, 1.0, [[(big, 40, c, True)]], align=PP_ALIGN.CENTER, space_after=0)
    text(s, cx+0.2, sy+1.5, sw-0.4, 0.8, [[(cap, 13, GREY, False)]],
         align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.1)

text(s, 0.7, 5.15, 11.9, 1.3,
     [[("“Managing highly nested, repetitive folder structures is prone to human "
        "error and consumes dozens of man-hours.”", 16, GREY, False)],
      [("SyncFlow Automator removes both — the structure and the naming become "
        "one-click, standardised, and repeatable.", 15, WHITE, True)]],
     space_after=8, line_spacing=1.1)


# =====================================================================
# SLIDE 11 — WHAT'S NEW v1.1
# =====================================================================
s = slide()
page_header(s, "Improvements", "What's new in v1.1 — safer & faster", 11)

improvements = [
    ("Extension safety", "Rename now preserves the real file extension — a .wav no "
     "longer becomes a mislabeled .mp3.", GREEN),
    ("Reliable capture", "Reads a file already on the clipboard first, then falls "
     "back to synthetic Ctrl+C — fixes wrong-window copies.", GREEN),
    ("Non-blocking search", "Tree is walked once on a background thread & filtered "
     "in memory; debounced keystrokes = no UI freeze.", BLUE),
    ("One-step Undo", "Instantly revert the last rename if a shot number was wrong.", BLUE),
    ("Settings persistence", "Last path / project / reel / characters are remembered "
     "between sessions.", AMBER),
    ("Safety + tests", "Large-deployment confirm prompt, overwrite guard, and a "
     "headless test suite (10/10 passing).", AMBER),
]
x = 0.7; y = 2.0; w = 5.9; h = 1.45; gx = 0.35; gy = 0.22
for i, (t, d, c) in enumerate(improvements):
    cx = x + (i % 2) * (w + gx)
    cy = y + (i // 2) * (h + gy)
    box(s, cx, cy, w, h, fill=PANEL, round_=True)
    box(s, cx, cy, 0.12, h, fill=c)
    text(s, cx+0.35, cy+0.16, w-0.6, 0.4, [[("✓  " + t, 15, WHITE, True)]], space_after=0)
    text(s, cx+0.35, cy+0.62, w-0.6, 0.8, [[(d, 12, GREY, False)]],
         space_after=0, line_spacing=1.03)


# =====================================================================
# SLIDE 12 — ROADMAP
# =====================================================================
s = slide()
page_header(s, "Roadmap", "Where it goes next", 12)

cols = [
    ("Now — shipped v1.1", GREEN, [
        "Extension-safe renaming",
        "Threaded quick search",
        "Undo + settings memory",
        "Headless test suite",
    ]),
    ("Next — high value", BLUE, [
        "Batch rename a whole folder",
        "Dry-run preview before apply",
        "Config file for formats/suffixes",
        "Custom naming templates",
    ]),
    ("Later — nice to have", AMBER, [
        "Drag-and-drop onto window",
        "Rename history / audit log",
        "Project presets per show",
        "Optional CLI mode",
    ]),
]
cw = 3.9; gx = 0.35; sx = 0.7; sy = 2.05; ch = 4.4
for i, (title, c, items) in enumerate(cols):
    cx = sx + i * (cw + gx)
    box(s, cx, sy, cw, ch, fill=PANEL, round_=True)
    box(s, cx, sy, cw, 0.6, fill=c, round_=True)
    text(s, cx, sy+0.08, cw, 0.5, [[(title, 14.5, WHITE, True)]],
         align=PP_ALIGN.CENTER, space_after=0)
    iy = sy + 0.9
    for it in items:
        box(s, cx+0.35, iy+0.09, 0.13, 0.13, fill=c)
        text(s, cx+0.62, iy-0.05, cw-0.9, 0.6, [[(it, 13, GREY, False)]],
             space_after=0, line_spacing=1.05)
        iy += 0.78


# =====================================================================
# SLIDE 13 — BUILD & DEPLOY
# =====================================================================
s = slide()
page_header(s, "Ship It", "Build & deploy — cross-platform", 13)

# Windows card
box(s, 0.7, 2.1, 5.9, 3.5, fill=PANEL, round_=True)
box(s, 0.7, 2.1, 5.9, 0.6, fill=BLUE, round_=True)
text(s, 0.7, 2.18, 5.9, 0.5, [[("🪟  Windows  (.exe)", 16, WHITE, True)]],
     align=PP_ALIGN.CENTER, space_after=0)
win_code = [
    "pip install pyinstaller",
    "",
    "pyinstaller --noconsole --onefile \\",
    '  --name=\"SyncFlow_Automator\" main.py',
    "",
    "→ dist/SyncFlow_Automator.exe",
]
text(s, 1.05, 3.0, 5.3, 2.4, [[(l, 12.5, GREEN if l.startswith("→") else GREY, False, MONO)]
                              for l in win_code], space_after=4, line_spacing=1.15)

# macOS card
box(s, 6.75, 2.1, 5.85, 3.5, fill=PANEL, round_=True)
box(s, 6.75, 2.1, 5.85, 0.6, fill=PANEL2, round_=True)
text(s, 6.75, 2.18, 5.85, 0.5, [[("  macOS  (.app / .dmg)", 16, WHITE, True)]],
     align=PP_ALIGN.CENTER, space_after=0)
mac_code = [
    "pyinstaller --noconsole --onefile \\",
    '  --name=\"SyncFlow_Automator\" main.py',
    "",
    "hdiutil create -format UDZO \\",
    "  -srcfolder dist/SyncFlow_Automator.app \\",
    "  dist/SyncFlow_Automator.dmg",
]
text(s, 7.1, 3.0, 5.2, 2.4, [[(l, 12, GREY, False, MONO)] for l in mac_code],
     space_after=4, line_spacing=1.15)

box(s, 0.7, 5.8, 11.9, 0.85, fill=PANEL, round_=True)
text(s, 1.0, 5.95, 11.4, 0.6,
     [[("Quick handover on Mac: no build needed — ", 13.5, GREY, False),
       ("python main.py", 13.5, BLUE_LT, True, MONO),
       ("  runs it straight from source (stdlib only).", 13.5, GREY, False)]],
     space_after=0)


# =====================================================================
# SLIDE 14 — CLOSE
# =====================================================================
s = slide()
accent_bar(s, h=0.2)
box(s, 0.9, 2.5, 1.4, 1.4, fill=BLUE, round_=True)
text(s, 0.9, 2.5, 1.4, 1.4, [[("⚡", 40, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, 2.55, 2.55, 10, 1.0, [[("Standardise the setup.", 34, WHITE, True)]], space_after=0)
text(s, 2.55, 3.35, 10, 1.0, [[("Ship the cut faster.", 34, BLUE_LT, True)]], space_after=0)

text(s, 2.6, 4.7, 10, 0.6,
     [[("One click builds the tree · one click names the file · one search finds "
        "anything.", 15, GREY, False)]], space_after=0)

chip(s, 2.6, 5.55, "Try:  python main.py", color=PANEL2, tw=3.4)
chip(s, 6.2, 5.55, "Tests:  python test_syncflow.py", color=PANEL2, tw=4.2)

text(s, 0.9, 6.75, 12, 0.4,
     [[("SyncFlow Automator v1.1.0  ·  Thank you", 12, GREY_DIM, False)]], space_after=0)


# ---------- save ----------
out = os.path.join(HERE, "SyncFlow_Automator_Deck.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
